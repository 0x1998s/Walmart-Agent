# 🛒 沃尔玛AI Agent平台 - 数据处理服务
# Walmart AI Agent Platform - Data Processing Service

import asyncio
import json
import logging
import mimetypes
from io import BytesIO
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple, Union
from uuid import uuid4

import pandas as pd
import numpy as np
from docx import Document as DocxDocument
from openpyxl import load_workbook
import PyPDF2
from pptx import Presentation

from app.core.config import get_settings
from app.services.vector_service import VectorService

settings = get_settings()
logger = logging.getLogger(__name__)


class DataService:
    """数据处理服务 - 处理结构化和非结构化数据"""
    
    def __init__(self):
        self.vector_service: Optional[VectorService] = None
        self.supported_formats = {
            # 文档格式
            '.pdf': self._process_pdf,
            '.docx': self._process_docx,
            '.doc': self._process_docx,  # 简化处理
            '.txt': self._process_text,
            '.md': self._process_text,
            
            # 表格格式
            '.xlsx': self._process_excel,
            '.xls': self._process_excel,
            '.csv': self._process_csv,
            
            # 演示文稿
            '.pptx': self._process_pptx,
            '.ppt': self._process_pptx,  # 简化处理
            
            # 数据格式
            '.json': self._process_json,
            '.xml': self._process_xml,
        }
    
    async def initialize(self, vector_service: VectorService):
        """初始化数据服务"""
        self.vector_service = vector_service
        logger.info("✅ 数据处理服务初始化完成")
    
    async def process_file(
        self,
        file_path: Union[str, Path],
        file_content: Optional[bytes] = None,
        metadata: Optional[Dict[str, Any]] = None,
        collection_name: str = "walmart_documents"
    ) -> Dict[str, Any]:
        """处理单个文件"""
        try:
            file_path = Path(file_path)
            file_extension = file_path.suffix.lower()
            
            if file_extension not in self.supported_formats:
                raise ValueError(f"不支持的文件格式: {file_extension}")
            
            # 读取文件内容
            if file_content is None:
                with open(file_path, 'rb') as f:
                    file_content = f.read()
            
            # 处理文件
            processor = self.supported_formats[file_extension]
            processed_data = await processor(file_content, file_path)
            
            # 准备元数据
            file_metadata = {
                "filename": file_path.name,
                "file_type": file_extension,
                "file_size": len(file_content),
                "processed_at": pd.Timestamp.now().isoformat(),
                **(metadata or {})
            }
            
            # 存储到向量数据库
            if self.vector_service and processed_data.get("text_chunks"):
                doc_ids = await self.vector_service.add_documents(
                    collection_name=collection_name,
                    documents=processed_data["text_chunks"],
                    metadatas=[
                        {**file_metadata, "chunk_index": i}
                        for i in range(len(processed_data["text_chunks"]))
                    ]
                )
                processed_data["document_ids"] = doc_ids
            
            logger.info(f"✅ 成功处理文件: {file_path.name}")
            return {
                "status": "success",
                "file_path": str(file_path),
                "metadata": file_metadata,
                **processed_data
            }
            
        except Exception as e:
            logger.error(f"❌ 处理文件失败 {file_path}: {e}")
            return {
                "status": "error",
                "file_path": str(file_path),
                "error": str(e)
            }
    
    async def process_batch_files(
        self,
        file_paths: List[Union[str, Path]],
        collection_name: str = "walmart_documents",
        batch_size: int = 10
    ) -> List[Dict[str, Any]]:
        """批量处理文件"""
        results = []
        
        for i in range(0, len(file_paths), batch_size):
            batch = file_paths[i:i + batch_size]
            
            # 并发处理批次
            batch_tasks = [
                self.process_file(file_path, collection_name=collection_name)
                for file_path in batch
            ]
            
            batch_results = await asyncio.gather(*batch_tasks, return_exceptions=True)
            
            for result in batch_results:
                if isinstance(result, Exception):
                    results.append({
                        "status": "error",
                        "error": str(result)
                    })
                else:
                    results.append(result)
        
        logger.info(f"✅ 批量处理完成，共处理 {len(file_paths)} 个文件")
        return results
    
    async def _process_pdf(self, content: bytes, file_path: Path) -> Dict[str, Any]:
        """处理PDF文件"""
        try:
            pdf_reader = PyPDF2.PdfReader(BytesIO(content))
            
            text_content = ""
            page_texts = []
            
            for page_num, page in enumerate(pdf_reader.pages):
                page_text = page.extract_text()
                text_content += page_text + "\n"
                page_texts.append({
                    "page": page_num + 1,
                    "text": page_text.strip()
                })
            
            # 分块处理
            text_chunks = self._chunk_text(text_content)
            
            return {
                "type": "pdf",
                "pages": len(pdf_reader.pages),
                "text_content": text_content.strip(),
                "text_chunks": text_chunks,
                "page_texts": page_texts,
                "metadata": {
                    "page_count": len(pdf_reader.pages),
                    "has_metadata": bool(pdf_reader.metadata),
                }
            }
            
        except Exception as e:
            raise RuntimeError(f"PDF处理失败: {e}")
    
    async def _process_docx(self, content: bytes, file_path: Path) -> Dict[str, Any]:
        """处理Word文档"""
        try:
            doc = DocxDocument(BytesIO(content))
            
            paragraphs = []
            text_content = ""
            
            for para in doc.paragraphs:
                if para.text.strip():
                    paragraphs.append(para.text.strip())
                    text_content += para.text + "\n"
            
            # 处理表格
            tables_data = []
            for table in doc.tables:
                table_data = []
                for row in table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    table_data.append(row_data)
                tables_data.append(table_data)
            
            text_chunks = self._chunk_text(text_content)
            
            return {
                "type": "docx",
                "text_content": text_content.strip(),
                "text_chunks": text_chunks,
                "paragraphs": paragraphs,
                "tables": tables_data,
                "metadata": {
                    "paragraph_count": len(paragraphs),
                    "table_count": len(tables_data),
                }
            }
            
        except Exception as e:
            raise RuntimeError(f"Word文档处理失败: {e}")
    
    async def _process_text(self, content: bytes, file_path: Path) -> Dict[str, Any]:
        """处理文本文件"""
        try:
            # 尝试不同编码
            encodings = ['utf-8', 'gbk', 'gb2312', 'latin-1']
            text_content = None
            
            for encoding in encodings:
                try:
                    text_content = content.decode(encoding)
                    break
                except UnicodeDecodeError:
                    continue
            
            if text_content is None:
                raise ValueError("无法解码文本文件")
            
            lines = text_content.splitlines()
            text_chunks = self._chunk_text(text_content)
            
            return {
                "type": "text",
                "text_content": text_content,
                "text_chunks": text_chunks,
                "lines": lines,
                "metadata": {
                    "line_count": len(lines),
                    "character_count": len(text_content),
                }
            }
            
        except Exception as e:
            raise RuntimeError(f"文本文件处理失败: {e}")
    
    async def _process_excel(self, content: bytes, file_path: Path) -> Dict[str, Any]:
        """处理Excel文件"""
        try:
            # 使用pandas读取Excel
            excel_data = pd.read_excel(BytesIO(content), sheet_name=None)
            
            sheets_data = {}
            text_content = ""
            
            for sheet_name, df in excel_data.items():
                # 转换为字典格式
                sheet_dict = df.to_dict('records')
                sheets_data[sheet_name] = sheet_dict
                
                # 提取文本内容
                sheet_text = f"工作表: {sheet_name}\n"
                sheet_text += df.to_string(index=False)
                text_content += sheet_text + "\n\n"
            
            text_chunks = self._chunk_text(text_content)
            
            return {
                "type": "excel",
                "text_content": text_content.strip(),
                "text_chunks": text_chunks,
                "sheets": sheets_data,
                "metadata": {
                    "sheet_count": len(excel_data),
                    "sheet_names": list(excel_data.keys()),
                }
            }
            
        except Exception as e:
            raise RuntimeError(f"Excel文件处理失败: {e}")
    
    async def _process_csv(self, content: bytes, file_path: Path) -> Dict[str, Any]:
        """处理CSV文件"""
        try:
            # 尝试不同编码
            encodings = ['utf-8', 'gbk', 'gb2312']
            df = None
            
            for encoding in encodings:
                try:
                    df = pd.read_csv(BytesIO(content), encoding=encoding)
                    break
                except UnicodeDecodeError:
                    continue
            
            if df is None:
                raise ValueError("无法解码CSV文件")
            
            # 转换为文本
            text_content = df.to_string(index=False)
            text_chunks = self._chunk_text(text_content)
            
            # 数据统计
            stats = {
                "row_count": len(df),
                "column_count": len(df.columns),
                "columns": df.columns.tolist(),
                "dtypes": df.dtypes.to_dict(),
            }
            
            return {
                "type": "csv",
                "text_content": text_content,
                "text_chunks": text_chunks,
                "data": df.to_dict('records'),
                "metadata": stats
            }
            
        except Exception as e:
            raise RuntimeError(f"CSV文件处理失败: {e}")
    
    async def _process_pptx(self, content: bytes, file_path: Path) -> Dict[str, Any]:
        """处理PowerPoint文件"""
        try:
            prs = Presentation(BytesIO(content))
            
            slides_data = []
            text_content = ""
            
            for slide_num, slide in enumerate(prs.slides):
                slide_text = f"幻灯片 {slide_num + 1}:\n"
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text += shape.text + "\n"
                
                slides_data.append({
                    "slide": slide_num + 1,
                    "text": slide_text.strip()
                })
                text_content += slide_text + "\n"
            
            text_chunks = self._chunk_text(text_content)
            
            return {
                "type": "pptx",
                "text_content": text_content.strip(),
                "text_chunks": text_chunks,
                "slides": slides_data,
                "metadata": {
                    "slide_count": len(prs.slides),
                }
            }
            
        except Exception as e:
            raise RuntimeError(f"PowerPoint文件处理失败: {e}")
    
    async def _process_json(self, content: bytes, file_path: Path) -> Dict[str, Any]:
        """处理JSON文件"""
        try:
            json_data = json.loads(content.decode('utf-8'))
            
            # 转换为文本表示
            text_content = json.dumps(json_data, ensure_ascii=False, indent=2)
            text_chunks = self._chunk_text(text_content)
            
            return {
                "type": "json",
                "text_content": text_content,
                "text_chunks": text_chunks,
                "data": json_data,
                "metadata": {
                    "json_structure": self._analyze_json_structure(json_data),
                }
            }
            
        except Exception as e:
            raise RuntimeError(f"JSON文件处理失败: {e}")
    
    async def _process_xml(self, content: bytes, file_path: Path) -> Dict[str, Any]:
        """处理XML文件"""
        try:
            import xml.etree.ElementTree as ET
            
            root = ET.fromstring(content.decode('utf-8'))
            
            # 提取文本内容
            text_content = self._extract_xml_text(root)
            text_chunks = self._chunk_text(text_content)
            
            return {
                "type": "xml",
                "text_content": text_content,
                "text_chunks": text_chunks,
                "metadata": {
                    "root_tag": root.tag,
                    "element_count": len(list(root.iter())),
                }
            }
            
        except Exception as e:
            raise RuntimeError(f"XML文件处理失败: {e}")
    
    def _chunk_text(self, text: str, chunk_size: int = 1000, overlap: int = 200) -> List[str]:
        """将文本分块"""
        if not text:
            return []
        
        chunks = []
        start = 0
        
        while start < len(text):
            end = start + chunk_size
            
            # 如果不是最后一块，尝试在句号处分割
            if end < len(text):
                last_period = text.rfind('.', start, end)
                last_newline = text.rfind('\n', start, end)
                split_point = max(last_period, last_newline)
                
                if split_point > start:
                    end = split_point + 1
            
            chunk = text[start:end].strip()
            if chunk:
                chunks.append(chunk)
            
            start = end - overlap
        
        return chunks
    
    def _analyze_json_structure(self, data: Any, max_depth: int = 3) -> Dict[str, Any]:
        """分析JSON结构"""
        if max_depth <= 0:
            return {"type": type(data).__name__}
        
        if isinstance(data, dict):
            return {
                "type": "dict",
                "keys": list(data.keys())[:10],  # 最多显示10个键
                "key_count": len(data),
                "sample_values": {
                    k: self._analyze_json_structure(v, max_depth - 1)
                    for k, v in list(data.items())[:3]  # 分析前3个值
                }
            }
        elif isinstance(data, list):
            return {
                "type": "list",
                "length": len(data),
                "sample_items": [
                    self._analyze_json_structure(item, max_depth - 1)
                    for item in data[:3]  # 分析前3个元素
                ] if data else []
            }
        else:
            return {"type": type(data).__name__, "value": str(data)[:100]}
    
    def _extract_xml_text(self, element) -> str:
        """从XML元素提取文本"""
        text_parts = []
        
        if element.text:
            text_parts.append(element.text.strip())
        
        for child in element:
            text_parts.append(self._extract_xml_text(child))
            if child.tail:
                text_parts.append(child.tail.strip())
        
        return ' '.join(filter(None, text_parts))
    
    async def search_documents(
        self,
        query: str,
        collection_name: str = "walmart_documents",
        filters: Optional[Dict[str, Any]] = None,
        n_results: int = 10
    ) -> Dict[str, Any]:
        """搜索文档"""
        if not self.vector_service:
            raise RuntimeError("向量服务未初始化")
        
        return await self.vector_service.hybrid_search(
            collection_name=collection_name,
            query=query,
            filters=filters,
            n_results=n_results
        )
    
    async def get_document_stats(self, collection_name: str = "walmart_documents") -> Dict[str, Any]:
        """获取文档统计信息"""
        if not self.vector_service:
            raise RuntimeError("向量服务未初始化")
        
        return await self.vector_service.get_collection_stats(collection_name)
